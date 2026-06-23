import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DATA_DIR = os.path.dirname(os.path.abspath(__file__))

def main():
    print("[START] Compiling Hackathon Presentation Slides...")
    prs = Presentation()
    
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ─────────────────────────────────────────────
    # HELPERS
    # ─────────────────────────────────────────────
    def add_blank_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Brownish white background color #f4ede2
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 237, 226)
        
        # Top banner shape in deep warm brown #4a3525
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(74, 53, 37)
        shape.line.color.rgb = RGBColor(50, 34, 22)
        
        # Title text box inside banner (white text for maximum contrast)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.33), Inches(0.66))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        return slide

    def add_bullet_column(slide, left, top, width, height, lines, font_size=13):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                
            p.font.name = 'Arial'
            p.font.size = Pt(font_size)
            
            # Format bullets
            if item.startswith("  - ") or item.startswith("  * "):
                p.level = 2
                p.text = item[4:]
                p.font.color.rgb = RGBColor(87, 96, 111) # medium charcoal #57606f
            elif item.startswith("- ") or item.startswith("* "):
                p.level = 1
                p.text = item[2:]
                p.font.color.rgb = RGBColor(47, 53, 66) # dark charcoal #2f3542
            else:
                p.level = 0
                p.text = item
                p.font.bold = True
                p.font.color.rgb = RGBColor(74, 53, 37) # deep warm brown headers
                p.space_before = Pt(8)
                
            # Apply color tags for highlights
            if "[ADVANCED FEATURE]" in p.text:
                p.text = p.text.replace("[ADVANCED FEATURE]", "")
                p.font.color.rgb = RGBColor(27, 94, 32) # deep forest green accent #1b5e20
                p.font.bold = True
                
        return txBox

    # ─────────────────────────────────────────────
    # SLIDE 1: TITLE SLIDE
    # ─────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(244, 237, 226)
    
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "GRIDLOCK-FREE BENGALURU"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Arial'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(74, 53, 37) # Deep warm brown
    
    p2 = tf.add_paragraph()
    p2.text = "Predictive Parking Intelligence, Anomaly Alerts & What-If Dispatch Simulator"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(47, 53, 66) # dark charcoal
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Bengaluru Traffic Police (BTP) Parking Violation Management Solution"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(87, 96, 111) # Slate grey
    p3.space_before = Pt(40)
    
    p4 = tf.add_paragraph()
    p4.text = "Solo Innovator | Dataset: 298,450 Records (Nov 2023 - Apr 2024) | Release: June 2026"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = 'Arial'
    p4.font.size = Pt(13)
    p4.font.color.rgb = RGBColor(27, 94, 32) # deep forest green
    p4.space_before = Pt(15)

    # ─────────────────────────────────────────────
    # SLIDE 2: THE PROBLEM & THE SCITA GAP
    # ─────────────────────────────────────────────
    slide2 = add_blank_slide("The Problem & The SCITA Enforcement Leakage")
    
    col1_lines = [
        "Systemic Enforcement Leakage",
        "- Raw violation analysis identified thousands of approved parking violations that were never sent to the SCITA challan system.",
        "- Leakage Formula: SCITA Gap = (validation_status == 'approved') & (data_sent_to_scita == False).",
        "- Consequence: Repeat parking offenders escape penalties, which directly undermines traffic rules and traffic flow.",
        "Illegal Parking Gridlock",
        "- Reactive parking management fails to prevent blockages at critical junctions.",
        "- Top zones include dense commercial sectors with high traffic volume (Shivajinagar, Upparpet, City Market)."
    ]
    add_bullet_column(slide2, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "Audit & Target Prioritization",
        "- Audit metrics pinpoint the exact station jurisdictions suffering from communication failures (SCITA Gap Count).",
        "- The pipeline prioritizes cells with large SCITA gaps to restore BTP's enforcement credibility.",
        "Manpower Allocation Solution",
        "- Rather than uniform deployment, BTP can target grid cells by priority to maximize enforcement coverage.",
        "- Coordinates patrols to shift windows (Morning, Evening, Night) based on local violation patterns."
    ]
    add_bullet_column(slide2, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SLIDE 3: 7-LAYER PIPELINE ARCHITECTURE
    # ─────────────────────────────────────────────
    slide3 = add_blank_slide("7-Layer Upgraded Pipeline Architecture")
    
    col1_lines = [
        "Layer 1: Preprocessing & Cleaning",
        "- Drops unused columns, converts date timestamps to Indian Standard Time (IST), and parses JSON violation types.",
        "Layer 2: Features & Spillover Index [ADVANCED FEATURE]",
        "- Establishes dual spatial grids: 3-decimal (~110m cells) for scoring and 2-decimal (~1.1km cells) for temporal forecast.",
        "- Computes junction multipliers and Congestion Spillover indices.",
        "Layer 3: Unsupervised & Supervised Models [ADVANCED FEATURE]",
        "- Runs DBSCAN spatial clustering, XGBoost temporal regressor, and Isolation Forest anomaly flagging.",
        "Layer 4: Composite Risk Scoring [ADVANCED FEATURE]",
        "- Safe MinMax normalization of density, severity, junctions, peak hours, and weights scaled by Congestion Spillover."
    ]
    add_bullet_column(slide3, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "Layer 5: Scheduler & Vehicle Asset Optimization [ADVANCED FEATURE]",
        "- Groups top 50 cells, pivots shifts, and assigns optimized patrol vehicles (Tows, Jeeps, Bikes) based on weight mode.",
        "Layer 6: Dynamic Dashboard with What-If Simulator [ADVANCED FEATURE]",
        "- Renders Map, Plotly charts, and unified dashboard incorporating dynamic JS re-ranking and Leaflet opacity filters.",
        "Layer 7: Outbox Email Dispatch Center",
        "- Simulates emails to police station heads, queued and dispatched 1 hour before shift, synced to timezone clock.",
        "Automated Verification Suite [ADVANCED FEATURE]",
        "- Runs verify_pipeline.py to assert schema fields, output file existence, and enforce design constraints."
    ]
    add_bullet_column(slide3, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SLIDE 4: DUAL MACHINE LEARNING LAYER
    # ─────────────────────────────────────────────
    slide4 = add_blank_slide("Unsupervised & Supervised Machine Learning Suite")
    
    col1_lines = [
        "Model A: Spatial Hotspot Clusterer (DBSCAN)",
        "- Filtered to approved records to find dense spatial violation zones.",
        "- Metric: Haversine distance on coordinates (eps=300m, min_samples=10).",
        "- Detected 175 major spatial hotspot clusters.",
        "- Performance Cache: DBSCAN labels are cached in dbscan_labels.npy to make dashboard load times instant on re-execution.",
        "Model B: Temporal Density Predictor (XGBoost)",
        "- Forecasts next-hour violation counts to serve as a traffic congestion proxy.",
        "- Inputs: Chronologically sorted dataset with 1h and 3h rolling lags computed using pandas transform.",
        "- Validation: Evaluated via TimeSeriesSplit (3 folds).",
        "- Metrics: Achieved validation RMSE of 26.08 and MAE of 11.06."
    ]
    add_bullet_column(slide4, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "Model C: Spatial-Temporal Anomaly Detector [ADVANCED FEATURE]",
        "- Algorithm: sklearn.ensemble.IsolationForest trained on cell-hour aggregates.",
        "- Training Features: Density, average severity, average vehicle weight, scita gap count, and peak-hour ratio.",
        "- Parameters: Fit with 5% contamination rate.",
        "- Flagged 786 anomalous cell-hour combinations.",
        "- Actionable Output: Mapped back to top 50 cells under is_anomaly. Anomalies blink distinctly on Leaflet map markers to catch operators' eyes immediately.",
        "- Alerts Feed: Dedicated sidebar panel lists anomalous zones in real-time."
    ]
    add_bullet_column(slide4, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SLIDE 5: ADVANCED INTELLIGENCE FEATURES
    # ─────────────────────────────────────────────
    slide5 = add_blank_slide("Advanced Intelligence & Optimization")
    
    col1_lines = [
        "Congestion Spillover Index [ADVANCED FEATURE]",
        "- Formulated to inflate risk scores of cells near major intersections to prevent traffic gridlock.",
        "- Algorithm: KDTree neighborhood search queries surrounding cells within a 300m bounding box (0.003 degrees).",
        "- Formula: Spillover Index = 1.0 + 0.15 * log1p(neighboring_junction_violations).",
        "- Scoring Fusion: Base Risk Score is multiplied by the Spillover Index, prioritizing zones adjacent to high-risk junctions.",
        "Vehicle-Type Patrol Optimization [ADVANCED FEATURE]",
        "- Classifies grid cells and optimizes vehicle assignments based on local vehicle weight distributions (mode).",
        "  - Heavy Tow Truck: For heavy mode grids (Tankers, Buses, Cabs).",
        "  - Patrol Jeep: For passenger grids (Cars, Autos).",
        "  - Interceptor Bike: For light grids (Motorcycles, Scooters)."
    ]
    add_bullet_column(slide5, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "Composite Risk Scoring Fusion",
        "- Evaluates cell priority using safe minmax normalization:",
        "  Risk = Base_Risk * Spillover_Index",
        "  Base_Risk = 40% Density + 25% Severity + 20% Junction + 10% Peak Hour + 5% Vehicle Weight",
        "Time-Synced Notifier System",
        "- Groups the schedule by station and shift to send targeted dispatch briefs.",
        "- Simulated email client tab in dashboard updates email status badges (Dispatched, Queued, Archived) dynamically based on the current timezone clock.",
        "- Emails contain detailed briefings: ranks, coordinates, expected hourly breakdown, and assigned patrol asset (Tow, Jeep, Bike)."
    ]
    add_bullet_column(slide5, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SLIDE 6: WHAT-IF SIMULATOR & DYNAMIC UI
    # ─────────────────────────────────────────────
    slide6 = add_blank_slide("What-If Dispatch Simulator & Front-End UI")
    
    col1_lines = [
        "Interactive Sliders [ADVANCED FEATURE]",
        "- Dispatchers can dynamically adjust available patrol units (5-50) and priority bias (Risk vs. Density) via sidebar sliders.",
        "- Score Blending: JavaScript blends normalized risk and density on slider change:",
        "  BlendedScore = Bias * Risk + (1 - Bias) * Density",
        "Violation Risk Mitigation Score (VRMS) [ADVANCED FEATURE]",
        "- Computes a real-time mitigation score in JS to measure the percentage of total risk covered by dispatched patrols:",
        "  VRMS = (Sum of Blended Scores of Dispatched Cells / Sum of Blended Scores of All 50 Cells) * 100",
        "- Rendered in a glowing neon progress bar that updates instantly."
    ]
    add_bullet_column(slide6, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "Leaflet Map Integration [ADVANCED FEATURE]",
        "- Dispatched patrol units remain at high opacity (0.8) and display rank-coded colors.",
        "- Standby units are dynamically faded to low opacity (0.15) to focus dispatchers' attention.",
        "Aesthetic UI/UX Enhancements [ADVANCED FEATURE]",
        "- Floating Map Overlay Switcher: Programmatically handles layer swapping (Static Heatmap vs Live Dispatch Heatmap) floated directly on the map card.",
        "- Anomaly Alerts Feed: Dedicated sidebar card shows live warnings.",
        "- Emoji Removal: Emojis are completely removed from all dashboard components and map controls for a professional look."
    ]
    add_bullet_column(slide6, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SLIDE 7: VERIFICATION & PACKAGE DELIVERABLES
    # ─────────────────────────────────────────────
    slide7 = add_blank_slide("Verification & Packaged Deliverables")
    
    col1_lines = [
        "Automated Verification Suite",
        "- Executing verify_pipeline.py runs the pipeline, checks input data, and runs structural asserts.",
        "- Schema Verification: Asserts that patrol_schedule.csv contains upgraded columns: is_anomaly, spillover_index, and assigned_patrol_unit.",
        "- Email Verification: Asserts that simulated_emails.log contains patrol vehicle assignments.",
        "- Output Checks: Asserts map.html, charts.html, dashboard.html, index.html, and patrol_schedule.pdf exist and are non-empty."
    ]
    add_bullet_column(slide7, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), col1_lines)
    
    col2_lines = [
        "BTP Enforcement Deliverables Archive",
        "- Deliverables are zipped into BTP_Enforcement_Pipeline.zip.",
        "- Design Bound Exclusions: Excludes large raw violations.csv (109MB), local index.html, and dbscan cache labels (dbscan_labels.npy) to ensure clean file transfer.",
        "- Packaged Files:",
        "  - run_pipeline.py & generate_pdf.py",
        "  - verify_pipeline.py",
        "  - dashboard.html, map.html & charts.html",
        "  - patrol_schedule.pdf & patrol_schedule.csv",
        "  - simulated_emails.log"
    ]
    add_bullet_column(slide7, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), col2_lines)

    # ─────────────────────────────────────────────
    # SAVE PRESENTATION
    # ─────────────────────────────────────────────
    out_path = os.path.join(DATA_DIR, "Gridlock_Free_Bengaluru.pptx")
    prs.save(out_path)
    print(f"[SUCCESS] PowerPoint slide deck created successfully: {out_path} ({os.path.getsize(out_path)/1024:.1f} KB)")

if __name__ == "__main__":
    main()
